"""
QA-1 — Gerador do Pitch ForwardService (v3 — 16 slides consecutivos por integrante)

Blocos por integrante (para edicao do video):
  Slides  1-4  : Joao Victor Franco — Problema + Contexto
  Slides  5-7  : Lucca Borges       — Inteligencia + ML
  Slides  8-11 : Ruan Melo          — Solucao + Arquitetura
  Slides 12-14 : Rodrigo Jimenez    — Valor + Canvas + Roadmap
  Slides 15-16 : Todos              — Equipe + Encerramento
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Cores ──────────────────────────────────────────────────────────────────────
FORD_DARK   = RGBColor(0x00, 0x27, 0x4A)
FORD_BLUE   = RGBColor(0x1E, 0x56, 0xA0)
FORD_ORANGE = RGBColor(0xF9, 0x63, 0x02)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
GRAY        = RGBColor(0x64, 0x74, 0x87)

# ── Equipe ─────────────────────────────────────────────────────────────────────
TEAM = [
    ("Rodrigo Jimenez",    "558148", "Testing/QA · Arquitetura TOGAF\nPitch · Canvas · Quadro de Valor\nCoordenação dos artefatos acadêmicos"),
    ("João Victor Franco", "556790", "Liderança técnica · Documentação\nSpecs TOGAF (suporte ao Rodrigo)\nIntegração entre disciplinas"),
    ("Lucca Borges",       "554608", "Machine Learning\nEDA · Feature Engineering\nK-means · XGBoost · Survival Analysis"),
    ("Ruan Melo",          "557599", "Backend Java\nSOA · REST + SOAP\nSpring Boot 3 · Flyway · Cyber"),
]

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh

def add_text(slide, text, l, t, w, h, size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def slide_header(slide, tag, label, presenter="", bg=FORD_DARK):
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
    add_rect(slide, 0, 0, 13.33, 1.0, bg)
    add_text(slide, tag,   0.3, 0.05, 8,    0.4,  size=11, bold=True, color=FORD_ORANGE)
    add_text(slide, label, 0.3, 0.45, 12,   0.5,  size=22, bold=True, color=WHITE)
    if presenter:
        add_text(slide, presenter, 10.5, 0.05, 2.7, 0.35, size=9, color=LIGHT_GRAY,
                 align=PP_ALIGN.RIGHT, italic=True)

def accent_line(slide, t=1.05):
    r = slide.shapes.add_shape(1, Inches(0), Inches(t), Inches(13.33), Inches(0.05))
    r.line.fill.background()
    r.fill.solid()
    r.fill.fore_color.rgb = FORD_ORANGE

def add_kpi_box(slide, value, label, l, t, w=2.8, h=1.3, bg=FORD_BLUE):
    add_rect(slide, l, t, w, h, bg)
    add_text(slide, value, l, t+0.05, w, 0.7,  size=30, bold=True,  color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t+0.75, w, 0.5,  size=11,             color=WHITE, align=PP_ALIGN.CENTER)

def add_bullets(slide, items, l, t, w, h, size=12, color=FORD_DARK):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "• " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCO JOAO VICTOR FRANCO (Slides 1–4)  — Abertura + Problema + Contexto
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 1 — Capa ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, FORD_DARK)
add_rect(s, 0, 5.8, 13.33, 1.7, FORD_BLUE)
add_rect(s, 0, 5.75, 13.33, 0.08, FORD_ORANGE)
add_text(s, "FORD CHALLENGE — FIAP 2026", 0.5, 0.4, 12, 0.5, size=13, color=FORD_ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "ForwardService", 0.5, 1.0, 12, 1.2, size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Plataforma de Retenção Pós-Venda Ford Brasil", 0.5, 2.2, 12, 0.7, size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Inteligência + Automação + ROI Mensurável", 0.5, 2.9, 12, 0.5, size=16, color=FORD_ORANGE, align=PP_ALIGN.CENTER, italic=True)
team_str = (
    f"{TEAM[0][0]} — RM: {TEAM[0][1]}          {TEAM[1][0]} — RM: {TEAM[1][1]}\n"
    f"{TEAM[2][0]} — RM: {TEAM[2][1]}             {TEAM[3][0]} — RM: {TEAM[3][1]}"
)
add_text(s, team_str, 0.5, 4.0, 12, 1.2, size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Testing, Compliance & Quality Assurance — Prof. Elias Bernardo", 0.5, 5.2, 12, 0.5, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Sprint 1 — Maio 2026     |     Vídeo Pitch: [INSERIR LINK APÓS GRAVAR]", 0.5, 6.05, 12, 0.5, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 2 — Problema Ford ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "PROBLEMA", "O problema exclusivo da Ford Brasil", "João Victor Franco")
accent_line(s)

add_kpi_box(s, "~12,4M", "veículos Ford\nem circulação",    0.4, 1.3, 2.7, 1.4, FORD_DARK)
add_kpi_box(s, "80%",    "são modelos\ndescontinuados",     3.3, 1.3, 2.7, 1.4, FORD_BLUE)
add_kpi_box(s, "145",    "concessionárias\n(era 283 em 2021)", 6.2, 1.3, 2.7, 1.4, FORD_DARK)
add_kpi_box(s, "~4%",    "VIN Share — só 4%\nda frota usa a rede", 9.1, 1.3, 2.7, 1.4, FORD_ORANGE)

add_text(s, "A Curva da Morte — retenção colapsa a cada revisão", 0.4, 2.9, 12, 0.45, size=15, bold=True, color=FORD_DARK)
add_bullets(s, [
    "1ª revisão: 31% dos VINs retornam   →   2ª: 22%   →   3ª: 15%   →   4ª: 9%",
    "Cada degrau perde 30–40% dos clientes anteriores  (dado real Ford, 602k eventos)",
    "Pós-garantia: retenção despenca de 82% para 20–40%  (Cox Automotive 2025)",
], 0.4, 3.35, 12.5, 1.5, size=14)
add_text(s, "A Ford é a única montadora no Brasil SEM programa de fidelização pré-pago.", 0.4, 5.05, 12.5, 0.5, size=15, bold=True, color=FORD_ORANGE)
add_text(s, "Cada VIN perdido = receita recorrente de R$ 400–2.000/revisão evaporando para oficinas independentes.", 0.4, 5.55, 12.5, 0.5, size=13, color=FORD_DARK)
add_text(s, "CAV: reter = R$ 1–5 | reconquistar = R$ 25–150 | adquirir novo = R$ 80–300", 0.4, 6.1, 12.5, 0.5, size=12, color=GRAY, italic=True)

# ── Slide 3 — Benchmarking ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "MERCADO", "Benchmarking — O que o mercado já faz", "João Victor Franco")
accent_line(s)
add_text(s, "O padrão-ouro existe — a Ford Brasil ainda não chegou lá", 0.4, 1.15, 12.5, 0.45, size=15, bold=True, color=FORD_DARK)

competitors = [
    ("NADA 20 Groups\n(EUA / Brasil via FENABRAVE)", [
        "Grupo de 20 dealers que se comparam em 200+ KPIs",
        "Já opera no Brasil — 52 concessionárias em 3 grupos",
        "Benchmarking financeiro: partes, mão de obra, absorção",
        "ForwardService adiciona: score de churn + ação automática",
    ]),
    ("Nissan Dealer Performance", [
        "30 KPIs com gamificação — espadas katana por categoria",
        "Dashboard em tempo real por dealer",
        "Recompensas progressivas por performance de retenção",
        "ForwardService vai além: predição, não só medição",
    ]),
    ("GM OnStar / Fidelidade Geral", [
        "Conectividade veicular com alertas proativos",
        "Recall + manutenção integrados ao app",
        "Previsão de serviço por telemetria",
        "Ford tem 80% da frota SEM telemetria — Fluxo Simplificado resolve",
    ]),
]
for (title, items), x in zip(competitors, [0.3, 4.55, 8.8]):
    add_rect(s, x, 1.7, 4.0, 5.4, FORD_BLUE)
    add_text(s, title, x+0.2, 1.78, 3.7, 0.8, size=13, bold=True, color=WHITE)
    add_rect(s, x, 2.55, 4.0, 0.04, FORD_ORANGE)
    add_bullets(s, items, x+0.2, 2.65, 3.7, 4.3, size=12, color=WHITE)

# ── Slide 4 — Diferenciações Competitivas ─────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "DIFERENCIAIS", "Por que a ForwardService é impossível de replicar com CRM genérico", "João Victor Franco")
accent_line(s)

logicas = [
    ("LN1 — Economia do VIN (LSV)", "LSV: Ka ~R$3.000 vs Ranger ~R$22.000.\nLeads priorizados por valor real, não volume."),
    ("LN2 — Curva da Morte",        "Disparo no momento crítico por perfil.\n4 perfis: Fiel · Econômico · Esquecido · Abandono."),
    ("LN3 — Rede Invertida",        "Mapeia desertos de serviço (VINs sem dealer).\n145 dealers vs 521 Fiat — a conta não fecha."),
    ("LN4 — Recall Gateway",        "3,4M recalls pendentes → porta de entrada.\n20–35% aceitam serviços adicionais."),
    ("LN5 — IHC (0–100)",           "Índice de Saúde da Concessionária:\nVIN Share 25% + Conversão 20% + NPS 15% + ..."),
    ("LN6 — Frota Descontinuada",   "80% da frota são Ka/Fiesta/EcoSport.\nSegmenta por fase com estratégia específica."),
    ("LN7 — Closed-Loop ROI",       "Cada ação rastreada do disparo ao resultado.\nR$1 investido → R$300+ por campanha WhatsApp."),
    ("LN8 — Flywheel de Dados",     "Mês 1: 70% → Mês 12: 85% → Ano 2: 90%.\nVantagem acumulativa a cada ciclo."),
    ("LN9 — Ponte Serviço-Venda",   "Clientes retidos recompram: 74% vs 44%.\nLTV de até US$175.000 por cliente."),
]
for i, (title, body) in enumerate(logicas):
    row, col = i // 3, i % 3
    x = 0.3 + col * 4.35
    y = 1.25 + row * 2.05
    bg = FORD_DARK if i % 2 == 0 else FORD_BLUE
    add_rect(s, x, y, 4.1, 1.85, bg)
    add_text(s, title, x+0.15, y+0.1,  3.85, 0.45, size=11, bold=True, color=FORD_ORANGE)
    add_text(s, body,  x+0.15, y+0.55, 3.85, 1.2,  size=10.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCO LUCCA BORGES (Slides 5–7)  — Inteligência + ML + Métricas
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 5 — Stack Técnica (foco Intelligence Hub / ML) ──────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "ARQUITETURA", "Stack Técnica + Intelligence Hub", "Lucca Borges")
accent_line(s)
add_text(s, "Princípio SOA: cada caixa tem responsabilidade única e fala só com o portão central (forward-api-java).", 0.4, 1.15, 12.5, 0.45, size=13, italic=True, color=FORD_DARK)

componentes = [
    ("forward-ml\n(Python · scikit-learn)", "INTELLIGENCE HUB",
     "K-means · XGBoost · Kaplan-Meier\n175k VINs × 20+ features comportamentais\nCurva da Morte: 31%→22%→15%→9% (dado real)"),
    ("forward-api-java\n(Java 17 · Spring Boot 3)", "PORTÃO CENTRAL",
     "REST + SOAP (SOA)\nJWT/RBAC · Flyway · SpringDoc\nÚnico ponto de entrada — ninguém fala direto com o banco"),
    ("forward-mobile\n(React Native · Expo)", "EXPERIENCE LAYER",
     "Modo Ford (atendente) / Modo Cliente (dono)\nPush notification · expo-router\nDual-mode — 1 app, 2 interfaces"),
    ("forward-web\n(SvelteKit)", "PERFORMANCE CONSOLE",
     "Cockpit Ford corporativo\nKPIs · Closed-Loop ROI · IHC 0-100\nRead-only · Chart.js / visx"),
    ("n8n · WhatsApp\n(self-hosted · Docker)", "ACTION ENGINE",
     "Orquestrador de ações automáticas\nWhatsApp Business · 97% abertura BR\nROI 300:1 por campanha"),
    ("PostgreSQL 16\n(Azure · Supabase/Neon)", "BANCO DE DADOS",
     "Fonte da verdade\naudit_log · vin_features\nSó o Java e o segmentador batch escrevem"),
]
for i, (name, role, desc) in enumerate(componentes):
    row, col = i // 3, i % 3
    x = 0.3 + col * 4.35
    y = 1.7 + row * 2.65
    bg = FORD_DARK if row == 0 else FORD_BLUE
    add_rect(s, x, y, 4.1, 2.45, bg)
    add_text(s, role, x+0.15, y+0.08, 3.85, 0.35, size=9,  bold=True, color=FORD_ORANGE)
    add_text(s, name, x+0.15, y+0.42, 3.85, 0.6,  size=13, bold=True, color=WHITE)
    add_text(s, desc, x+0.15, y+1.05, 3.85, 1.3,  size=10, color=WHITE)

# ── Slide 6 — Critérios de Qualidade ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "QUALIDADE", "Critérios de Qualidade — Performance · Segurança · Usabilidade", "Lucca Borges")
accent_line(s)

areas = [
    ("PERFORMANCE", FORD_DARK, [
        ("Latência API p95",         "< 300 ms",   "Medido em produção Azure"),
        ("Disponibilidade",          "≥ 99%",       "Docker + health checks"),
        ("Throughput segmentador",   "175k VINs/dia","Cron 02h — batch diário"),
        ("Falsos positivos churn",   "< 10%",       "Threshold tunable por perfil"),
        ("ML accuracy (AUC)",        "0,82–0,90",   "XGBoost + SHAP interpretável"),
    ]),
    ("SEGURANÇA (CYBER)", FORD_BLUE, [
        ("Autenticação",             "JWT HS256",    "RBAC por role + dealer"),
        ("Dados pessoais",           "Pseudonimizado","VIN_Hash SHA1 + salt Ford"),
        ("Comunicação",              "TLS 1.2+",     "HTTPS forçado em toda a infra"),
        ("Auditoria",                "audit_log",    "Trilha de cada acesso/ação"),
        ("Rate limiting",            "Bucket4j",     "Por IP + user — anti-DDoS"),
    ]),
    ("USABILIDADE", RGBColor(0x0D, 0x3B, 0x6E), [
        ("NPS alvo (cliente)",       "> 50",        "Pesquisa pós-atendimento"),
        ("Fluxo simplificado",       "Ka/EcoSport", "80% da frota sem telemetria"),
        ("Resposta WhatsApp",        "< 30s",       "n8n automático 24/7"),
        ("Onboarding atendente",     "< 15 min",    "UX dual-mode intuitivo"),
        ("i18n",                     "PT-BR + EN",  "Expo i18n desde o dia 1"),
    ]),
]
for i, (title, bg, metrics) in enumerate(areas):
    x = 0.3 + i * 4.35
    add_rect(s, x, 1.25, 4.1, 5.9, bg)
    add_text(s, title, x+0.2, 1.35, 3.8, 0.5, size=14, bold=True, color=FORD_ORANGE)
    add_rect(s, x, 1.85, 4.1, 0.04, WHITE)
    for j, (metric, value, note) in enumerate(metrics):
        my = 2.0 + j * 0.85
        add_text(s, metric, x+0.15, my+0.02, 2.3,  0.35, size=10, color=WHITE)
        add_text(s, value,  x+2.45, my+0.02, 1.4,  0.35, size=10, bold=True, color=FORD_ORANGE)
        add_text(s, note,   x+0.15, my+0.38, 3.8,  0.3,  size=9, color=LIGHT_GRAY, italic=True)

# ── Slide 7 — Métricas de Sucesso ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "MÉTRICAS", "Métricas de Sucesso — O que medimos", "Lucca Borges")
accent_line(s)
add_text(s, "ForwardService é mensurável por design — cada real investido tem retorno rastreável (Closed-Loop ROI)", 0.4, 1.15, 12.5, 0.45, size=13, italic=True, color=FORD_DARK)

metricas = [
    ("RETENÇÃO", [
        "VIN Share: 3,5–5% → meta 8–10% em 24 meses",
        "Taxa de retorno pós-1ª revisão: 71% → meta 80%",
        "+5% retenção = +25% a 95% de lucro (Bain)",
        "Contratos Ford Care: retenção 3× (20% → 60%)",
    ]),
    ("NPS & SATISFAÇÃO", [
        "NPS alvo > 50 (benchmark: Tesla 96, Toyota 60)",
        "Tempo de resposta WhatsApp < 30 segundos 24/7",
        "Conversão campanha WhatsApp: 30–45%",
        "Redução de falsos positivos < 10% por perfil",
    ]),
    ("ROI FINANCEIRO", [
        "ROI por campanha WhatsApp: 300:1 ou mais",
        "CAV: R$ 1–5 (ativo) vs R$ 25–150 (reconquista)",
        "Ticket médio recall convertido: R$ 750–1.500 extra",
        "49% do lucro bruto dealer vem do pós-venda (NADA)",
    ]),
    ("ML / QUALIDADE", [
        "AUC classificador: 0,82–0,90 (XGBoost+SHAP)",
        "Silhouette K-means > 0,4 (k=4 clusters)",
        "Precisão mês 1: 70% → mês 12: 85% → ano 2: 90%",
        "Latência API p95 < 300ms · disponibilidade ≥ 99%",
    ]),
]
for i, (title, items) in enumerate(metricas):
    row, col = i // 2, i % 2
    x = 0.3 + col * 6.5
    y = 1.7 + row * 2.7
    bg = FORD_DARK if i % 2 == 0 else FORD_BLUE
    add_rect(s, x, y, 6.25, 2.5, bg)
    add_text(s, title, x+0.2, y+0.1, 5.9, 0.45, size=13, bold=True, color=FORD_ORANGE)
    add_rect(s, x, y+0.55, 6.25, 0.04, WHITE)
    add_bullets(s, items, x+0.2, y+0.65, 5.9, 1.75, size=11.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCO RUAN MELO (Slides 8–11)  — Solução + Narrativa + TOGAF + Riscos
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 8 — Solução (4 Pilares) ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "SOLUÇÃO", "ForwardService — Os 4 Pilares", "Ruan Melo")
accent_line(s)

pilares = [
    ("INTELLIGENCE HUB",  "Quem está prestes a sair?\nK-means + XGBoost + Survival Analysis\nsobre 175k VINs reais Ford BR"),
    ("ACTION ENGINE",      "O que fazer, agora, automaticamente?\nWhatsApp · n8n · Pulse Leads\nROI 300:1 por campanha"),
    ("EXPERIENCE LAYER",   "Por que o cliente vai querer voltar?\nFord Care (plano pré-pago)\nFluxo Simplificado (Ka/Fiesta/EcoSport)"),
    ("PERFORMANCE CONSOLE","Funcionou? Quanto retornou?\nClosed-Loop ROI · IHC 0-100\nFlywheel — fica mais inteligente a cada ciclo"),
]
for (title, body), x in zip(pilares, [0.3, 3.6, 6.85, 10.1]):
    add_rect(s, x, 1.25, 2.9, 5.5, FORD_DARK)
    add_text(s, title, x+0.15, 1.35, 2.6, 0.7, size=11, bold=True, color=FORD_ORANGE)
    add_rect(s, x, 2.05, 2.9, 0.04, FORD_BLUE)
    add_text(s, body,  x+0.15, 2.15, 2.65, 4.5, size=11.5, color=WHITE)

add_text(s, "O fluxo é cíclico — cada ciclo alimenta o próximo com dados melhores (Flywheel).", 0.4, 7.0, 12.5, 0.45, size=12, italic=True, color=FORD_DARK, align=PP_ALIGN.CENTER)

# ── Slide 9 — Narrativa (Um dia no ForwardService) ────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "NARRATIVA", "Um dia no ForwardService — o produto 100% rodando", "Ruan Melo")
accent_line(s)

personas = [
    ("João Silva\nKa 2014 · sumido há 11 meses", FORD_DARK,
     "02h — Sistema re-etiqueta João: Esquecido, risco 78.\n"
     "09h12 — WhatsApp automático: '20% off na revisão até sexta.'\n"
     "09h47 — João responde. Termômetro cai de 78 para 71.\n"
     "Sexta — Revisão feita. Etiqueta muda para Fiel. Score: 22.\n"
     "O ciclo fechou. O dado alimenta o cérebro."),
    ("Maria\nAtendente · Concessionária Vila Olímpia", FORD_BLUE,
     "11h05 — App mobile, Modo Ford.\n"
     "'João — Esquecido — risco 71 — começou a responder no WhatsApp.'\n"
     "Maria liga, fecha agendamento pra sexta.\n"
     "Marca 'agendado' no app. Sistema atualizado.\n"
     "Uma ação, muitos sistemas informados."),
    ("Carlos\nDiretor Regional Ford", RGBColor(0x0D, 0x3B, 0x6E),
     "18h00 — Painel web, Modo Ford Corporativo.\n"
     "'47 esquecidos contatados · 22 agendaram · 47% conversão'\n"
     "'R$ 184 mil em revisão prevista'\n"
     "'9 dealers acima da meta · 3 abaixo'\n"
     "Carlos vê a frota inteira, não o João individual."),
]
for i, (name, bg, story) in enumerate(personas):
    x = 0.3 + i * 4.35
    add_rect(s, x, 1.25, 4.1, 5.8, bg)
    add_rect(s, x, 1.25, 4.1, 0.06, FORD_ORANGE)
    add_text(s, name,  x+0.2, 1.35, 3.8, 0.75, size=13, bold=True, color=WHITE)
    add_rect(s, x, 2.1, 4.1, 0.04, WHITE)
    add_text(s, story, x+0.2, 2.22, 3.8, 4.7,  size=11.5, color=WHITE)

add_text(s, "\"Uma porta, dois lados.  Um app, infinitas conversas.\"", 0.5, 7.0, 12.3, 0.45, size=13, italic=True, color=FORD_DARK, align=PP_ALIGN.CENTER)

# ── Slide 10 — Arquitetura TOGAF (6 Vistas) ───────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "TOGAF / ARCHI", "Arquitetura da Solução — 6 Vistas ArchiMate 3.1", "Ruan Melo")
accent_line(s)
add_text(s, "Visão · Negócio · Aplicação · Tecnologia · Requisitos de Qualidade · Monitoramento", 0.4, 1.12, 12.5, 0.35, size=12, italic=True, color=FORD_DARK, align=PP_ALIGN.CENTER)

togaf_views = [
    ("01 — Motivation\n(Drivers, Stakeholders, Goals)",
     r"C:\Users\USER\Desktop\espm\forward-docs\academic\togaf\views\01_motivation_view.png"),
    ("02 — Business\n(Processo As-Is vs To-Be)",
     r"C:\Users\USER\Desktop\espm\forward-docs\academic\togaf\views\02_business_view.png"),
    ("03 — Application\n(Componentes + Comunicações)",
     r"C:\Users\USER\Desktop\espm\forward-docs\academic\togaf\views\03_application_view.png"),
    ("04 — Technology\n(Infra: Azure, Docker, TLS)",
     r"C:\Users\USER\Desktop\espm\forward-docs\academic\togaf\views\04_technology_view.png"),
    ("05 — Requirements\n(Qualidade → Componentes)",
     r"C:\Users\USER\Desktop\espm\forward-docs\academic\togaf\views\05_requirements_view.png"),
    ("06 — Monitoring\n(Observabilidade em Produção)",
     r"C:\Users\USER\Desktop\espm\forward-docs\academic\togaf\views\06_monitoring_view.png"),
]
img_w, img_h = 4.1, 2.3
for i, (label, img_path) in enumerate(togaf_views):
    row, col = i // 3, i % 3
    x = 0.1 + col * (img_w + 0.28)
    y = 1.55 + row * (img_h + 0.65)
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(img_w), Inches(img_h))
    else:
        add_rect(s, x, y, img_w, img_h, FORD_BLUE)
        add_text(s, "[PNG não encontrado]", x, y+0.8, img_w, 0.5, size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, y + img_h + 0.03, img_w, 0.52, FORD_DARK)
    add_text(s, label, x+0.08, y + img_h + 0.05, img_w - 0.12, 0.46, size=8, color=WHITE)

add_text(s, "Arquivo editável: ForwardService.archimate (abre no Archi 4.x, gratuito em archimatetool.com)", 0.3, 7.08, 12.7, 0.35, size=9, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ── Slide 11 — Riscos e Mitigações ───────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "RISCOS", "Riscos da Solução + Mitigações", "Ruan Melo")
accent_line(s)

riscos = [
    ("Dataset sem labels prontos", "ALTO", FORD_ORANGE,
     "Target engenheirado (days_since_last_service > 365) + clustering cego. Documentado honestamente no MODEL_CARD."),
    ("Aprovação de template WhatsApp Meta", "MÉDIO", FORD_BLUE,
     "Sandbox Meta para demo. Templates aprovados previamente. Screenshot de pendência como fallback honesto."),
    ("DealerCode sem mapeamento", "MÉDIO", FORD_BLUE,
     "Agrupamento por UF via heurística. 80 dealers geocodificados (ABRADIF). Mapeamento exato no Sprint 2."),
    ("Ferramenta Archi nova", "BAIXO", FORD_DARK,
     "Export PNG como fallback. XML do modelo entregue. Suporte do João Victor com specs MD das 4 views."),
    ("Retenção baixa mesmo com a plataforma", "BAIXO", FORD_DARK,
     "Planos Ford Care elevam retenção de 20% para 60% (3x). Fluxo Simplificado alcança 80% da frota."),
    ("Segurança LGPD", "MÉDIO", FORD_BLUE,
     "VIN_Hash SHA1 robusto (5M reversões = 0 match). audit_log completo. Política descrita no plano Cyber."),
]
for i, (risk, level, bg, mit) in enumerate(riscos):
    row, col = i // 2, i % 2
    x = 0.3 + col * 6.5
    y = 1.25 + row * 2.1
    add_rect(s, x, y, 6.25, 1.9, bg)
    add_text(s, f"[{level}]  {risk}", x+0.15, y+0.1, 5.8, 0.55, size=12, bold=True, color=WHITE)
    add_rect(s, x, y+0.65, 6.25, 0.04, FORD_ORANGE)
    add_text(s, mit, x+0.15, y+0.75, 5.9, 1.0, size=10.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCO RODRIGO JIMENEZ (Slides 12–14)  — Valor + Canvas + Roadmap
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 12 — Business Model Canvas (resumo) ─────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(s, 0, 0, 13.33, 0.85, FORD_DARK)
add_text(s, "CANVAS  |  Rodrigo Jimenez", 0.3, 0.05, 8, 0.35, size=10, bold=True, color=FORD_ORANGE)
add_text(s, "Business Model Canvas — ForwardService", 0.3, 0.42, 12, 0.38, size=18, bold=True, color=WHITE)
add_rect(s, 0, 0.85, 13.33, 0.04, FORD_ORANGE)

def cblock(slide, title, items, x, t, w, h, bg=FORD_DARK):
    add_rect(slide, x, t, w, h, bg)
    add_text(slide, title, x+0.1, t+0.05, w-0.15, 0.28, size=8, bold=True, color=FORD_ORANGE)
    add_rect(slide, x, t+0.33, w, 0.02, FORD_ORANGE if bg != FORD_BLUE else WHITE)
    tb = slide.shapes.add_textbox(Inches(x+0.1), Inches(t+0.36), Inches(w-0.15), Inches(h-0.42))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "• " + it
        p.font.size = Pt(8)
        p.font.color.rgb = WHITE
        p.space_after = Pt(1)

cblock(s, "PARCEIROS-CHAVE", [
    "Ford Brasil (dataset oficial)", "Meta — WhatsApp Business (97%)",
    "Microsoft Azure (infra VM)", "Supabase / Neon (Postgres 16)",
    "FENABRAVE (dados frota BR)", "OpenAI / Anthropic (futuro)",
], 0.1, 0.92, 2.1, 3.3, FORD_DARK)

cblock(s, "ATIVIDADES-CHAVE", [
    "ML scoring 175k VINs/dia", "Disparos n8n + WhatsApp auto",
    "Manter e retreinar modelos (Flywheel)", "KPIs e IHC por dealer real time",
], 2.28, 0.92, 2.1, 1.57, FORD_BLUE)

cblock(s, "RECURSOS-CHAVE", [
    "Dataset 602k eventos Ford", "Modelos ML (3 camadas)",
    "Infra Azure VM + Postgres", "Equipe 4 devs multidisciplinar",
], 2.28, 2.57, 2.1, 1.65, FORD_DARK)

add_rect(s, 4.46, 0.92, 2.55, 3.3, FORD_BLUE)
add_rect(s, 4.46, 0.92, 2.55, 0.06, FORD_ORANGE)
add_text(s, "PROPOSTA DE VALOR", 4.56, 0.97, 2.4, 0.28, size=8, bold=True, color=FORD_ORANGE)
add_rect(s, 4.46, 1.25, 2.55, 0.02, WHITE)
vp = ["Ford BR: VIN Share 4% → 8-10%", "Dealer: Pulse Leads diários por risco",
      "Cliente Ka/EcoSport: Fluxo Simplificado", "Cliente Ranger: Ford Care (3× retenção)",
      "Cockpit nacional em tempo real", "ROI 300:1 por campanha WhatsApp"]
tb = s.shapes.add_textbox(Inches(4.56), Inches(1.3), Inches(2.38), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
first = True
for it in vp:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False; p.text = "• " + it
    p.font.size = Pt(8); p.font.bold = True
    p.font.color.rgb = WHITE; p.space_after = Pt(2)

cblock(s, "RELACIONAMENTO", [
    "Self-service: app dual-mode", "Personalizado via dealer",
    "WhatsApp automático 24/7", "Ford Care (longo prazo pré-pago)",
], 7.09, 0.92, 2.1, 1.57, FORD_DARK)

cblock(s, "CANAIS", [
    "WhatsApp Business (97% abertura)", "App Mobile React Native/Expo",
    "Dashboard Web SvelteKit", "n8n (email/SMS — Sprint 2)",
], 7.09, 2.57, 2.1, 1.65, FORD_BLUE)

cblock(s, "SEGMENTOS", [
    "Ford Brasil (B2B — matriz)", "145 concessionárias ativas",
    "Donos modelos recentes (Ranger)", "Donos Ka/EcoSport/Fiesta (2,5M+)",
    "Gestores regionais Ford", "4 perfis ML: Fiel/Econom/Esquec/Aband",
], 9.27, 0.92, 3.96, 3.3, FORD_DARK)

cblock(s, "ESTRUTURA DE CUSTOS", [
    "Azure VM: $7-30/mês · Supabase: $0-25/mês",
    "WhatsApp API: free 1k msgs/mês (Sprint 1)",
    "Retrabalho ML, suporte dealers, incidentes",
    "Treinamento atendentes: <15 min/usuário",
], 0.1, 4.3, 6.6, 2.0, FORD_BLUE)

cblock(s, "FONTES DE RECEITA", [
    "SaaS Ford Brasil: licença mensal pós-MVP",
    "Assinatura por dealer: 145 potenciais",
    "Ford Care: receita antecipada de revisões",
    "ROI 300:1 documentado por campanha",
], 6.78, 4.3, 6.45, 2.0, FORD_DARK)

add_text(s, "Detalhes completos (impacto de falhas + custos de qualidade): CANVAS.html / CANVAS.pdf", 0.2, 6.38, 12.9, 0.3, size=8, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ── Slide 13 — Quadro de Valor ────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "QUADRO DE VALOR", "Partes Interessadas — O que cada stakeholder ganha", "Rodrigo Jimenez")
accent_line(s)
add_text(s, "Para cada benefício: métrica de negócio (resultado) + métrica de qualidade (SLA técnico)", 0.4, 1.1, 12.5, 0.35, size=12, italic=True, color=FORD_DARK)

hdr_y = 1.5
cols = [("STAKEHOLDER", 0.15, 2.1), ("PROMESSA FORWARDSERVICE", 2.3, 3.5),
        ("MÉTRICA DE NEGÓCIO", 5.85, 3.2), ("MÉTRICA DE QUALIDADE", 9.1, 2.8), ("PRIO", 11.95, 1.2)]
for label, x, w in cols:
    add_rect(s, x, hdr_y, w, 0.38, FORD_DARK)
    add_text(s, label, x+0.05, hdr_y+0.04, w-0.1, 0.3, size=8, bold=True, color=WHITE)

rows_qv = [
    ("Ford Brasil",          "Dashboard nacional real time:\nVIN Share por UF · ROI rastreável",
     "VIN Share ≥ 8% em 24 meses\n(base: ~3,5-5% hoje)", "Disponib. ≥ 99%\nLatência API p95 < 300ms", "ALTA", FORD_ORANGE),
    ("Gestor Regional",      "Cockpit: IHC 0-100 por dealer,\nranking e causa-raiz",
     "≥ 80% dealers acima\nda meta regional de IHC", "IHC atualiz. 1×/dia\nSilhouette K-means > 0,4", "ALTA", FORD_ORANGE),
    ("Dono de Concessionária","Pulse Leads: lista diária\npor risco + LSV automático",
     "Conversão lead → revisão\n≥ 35%", "Endpoint /leads < 200ms\nZero leads duplicados/dia", "ALTA", FORD_ORANGE),
    ("Gerente de Serviço",   "Ficha do cliente: etiqueta de\nsegmento + ação sugerida pelo ML",
     "Redução 40% no tempo\nde priorização de contatos", "AUC classificador ≥ 0,82\nFalsos positivos < 10%", "ALTA", FORD_ORANGE),
    ("Dono modelo descontin.","Fluxo Simplificado: WhatsApp,\ncadastro VIN manual, lembrete km",
     "Reativação ≥ 15% dos VINs\ninativos >12 meses (1º sem.)", "Entrega WhatsApp < 30s\nAbertura ≥ 80%", "ALTA", FORD_ORANGE),
    ("Atendente (dealer)",   "Vista 360: histórico completo,\nscore, ação — uma tela",
     "NPS do atendimento ≥ 60\n(benchmark setor: ~50)", "Vista 360 < 1s\nOnboarding < 15 min", "MÉDIA", FORD_BLUE),
    ("Dono modelo recente",  "App: revisão por km/meses,\nFord Care pré-pago, agendam. 1 toque",
     "Adoção Ford Care ≥ 20%\ndos elegíveis em 12 meses", "App carrega < 2s\nPush notif. < 5s", "MÉDIA", FORD_BLUE),
    ("Equipe ForwardService", "Flywheel: precisão cresce com uso.\nSOA: cada componente evolui sozinho",
     "Precisão ML: 70%→80%→90%\n(mês 1 → 6 → ano 2)", "Cobertura testes ≥ 70%\nDeploy zero downtime", "MÉDIA", FORD_BLUE),
]
row_h = 0.61
for i, (stk, promessa, met_neg, met_qual, prio, prio_color) in enumerate(rows_qv):
    ry = hdr_y + 0.4 + i * row_h
    bg_row = RGBColor(0xF7, 0xFA, 0xFC) if i % 2 == 0 else WHITE
    add_rect(s, 0.15, ry, 13.0, row_h - 0.02, bg_row)
    add_text(s, stk,       0.2,  ry+0.03, 2.0,  row_h-0.08, size=8.5, bold=True, color=FORD_DARK)
    add_text(s, promessa,  2.35, ry+0.03, 3.35, row_h-0.08, size=8,             color=FORD_DARK)
    add_text(s, met_neg,   5.9,  ry+0.03, 3.1,  row_h-0.08, size=8,  bold=True, color=FORD_BLUE)
    add_text(s, met_qual,  9.15, ry+0.03, 2.65, row_h-0.08, size=8,  bold=True, color=RGBColor(0x27, 0x67, 0x49))
    add_rect(s, 12.0, ry+0.08, 1.1, row_h-0.22, prio_color)
    add_text(s, prio, 12.0, ry+0.1, 1.1, row_h-0.25, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Documento completo: QUADRO_DE_VALOR.html / PDF", 0.3, 6.55, 12.7, 0.3, size=8, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ── Slide 14 — Roadmap ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "ROADMAP", "Roadmap — Sprint 2, 3 e 4", "Rodrigo Jimenez")
accent_line(s)

fases = [
    ("SPRINT 1\n24/05/2026 (atual)", FORD_ORANGE, [
        "ML completo: EDA + K-means + XGBoost + Survival",
        "Backend Java: REST + SOAP + Flyway + Swagger",
        "Mobile: 3 telas dual-mode (Ford + Cliente)",
        "Dashboard web: Visão da frota",
        "n8n: 1 flow WhatsApp live",
        "Infra: Azure VM + Postgres",
        "Pitch + Canvas + TOGAF + Vídeo (QA)",
    ]),
    ("SPRINT 2\nJunho 2026", FORD_BLUE, [
        "Recomendador via ML (substitui regras)",
        "Performance Console completo (drill-down, export)",
        "Recall Gateway com 12 campanhas Senacon",
        "Multi-canal: SMS + email + push",
        "Closed-Loop ROI completo",
        "Anonimização full (LGPD) em repouso",
        "Testes automatizados end-to-end",
    ]),
    ("SPRINT 3–4\nJulho–Agosto 2026", FORD_DARK, [
        "Integração com DMS (quando Ford disponibilizar)",
        "Rede Invertida + mapeamento de desertos",
        "Ford Care (planos pré-pagos) no app",
        "Ponte Serviço-Venda (lead de recompra)",
        "Admin de configuração de regras",
        "MFA + criptografia avançada",
        "Piloto real com dealer parceiro",
    ]),
]
for i, (title, bg, items) in enumerate(fases):
    x = 0.3 + i * 4.35
    add_rect(s, x, 1.25, 4.1, 5.9, bg)
    add_text(s, title, x+0.2, 1.35, 3.8, 0.75, size=13, bold=True, color=WHITE)
    add_rect(s, x, 2.1, 4.1, 0.04, WHITE)
    tb = s.shapes.add_textbox(Inches(x+0.2), Inches(2.2), Inches(3.85), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.text = "✓  " + it
        p.font.size = Pt(11.5); p.font.color.rgb = WHITE; p.space_after = Pt(5)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCO TODOS (Slides 15–16)  — Equipe + Encerramento
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 15 — Equipe ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
slide_header(s, "EQUIPE", "Equipe ForwardService — FIAP 2026")
accent_line(s)

for (name, rm, role), x in zip(TEAM, [0.4, 3.5, 6.6, 9.7]):
    add_rect(s, x, 1.3, 3.1, 5.5, FORD_DARK)
    add_rect(s, x, 1.3, 3.1, 0.06, FORD_ORANGE)
    add_rect(s, x+0.55, 1.5, 2.0, 2.0, FORD_BLUE)
    add_text(s, "👤", x+0.55, 1.7, 2.0, 1.5, size=36, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, name,       x+0.15, 3.6,  2.8, 0.7, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, f"RM: {rm}", x+0.15, 4.3,  2.8, 0.4, size=11, color=FORD_ORANGE, align=PP_ALIGN.CENTER)
    add_rect(s, x, 4.72, 3.1, 0.04, FORD_BLUE)
    add_text(s, role,       x+0.15, 4.82, 2.8, 1.8, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "Disciplina: Testing, Compliance & Quality Assurance — Prof. Elias Bernardo", 0.4, 7.0, 12.5, 0.45, size=12, color=FORD_DARK, italic=True, align=PP_ALIGN.CENTER)

# ── Slide 16 — Encerramento ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, FORD_DARK)
add_rect(s, 0, 0, 13.33, 1.0, FORD_BLUE)
add_rect(s, 0, 0.95, 13.33, 0.07, FORD_ORANGE)
add_text(s, "PRÓXIMOS PASSOS", 0.4, 0.05, 12.5, 0.45, size=11, bold=True, color=FORD_ORANGE)
add_text(s, "O que vem depois do Sprint 1", 0.4, 0.45, 12.5, 0.5, size=22, bold=True, color=WHITE)

proximos = [
    "Sprint 2 — Recomendador via ML + Recall Gateway com 12 campanhas + Performance Console completo",
    "Sprint 3 — Integração DMS (quando Ford disponibilizar) + Ford Care no app + Rede Invertida",
    "Sprint 4 — Piloto real com dealer parceiro + MFA + criptografia avançada + Ponte Serviço-Venda",
    "Visão futura — FastAPI para o Classificador · Multi-idioma EN/ES · Admin de regras · Dark mode",
]
for i, item in enumerate(proximos):
    y = 1.2 + i * 0.9
    add_rect(s, 0.4, y, 0.45, 0.45, FORD_ORANGE)
    add_text(s, str(i+1), 0.4, y, 0.45, 0.45, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, item, 1.0, y+0.02, 12.0, 0.75, size=13, color=WHITE)

add_rect(s, 0, 4.9, 13.33, 2.6, FORD_BLUE)
add_rect(s, 0, 4.87, 13.33, 0.06, FORD_ORANGE)
add_text(s, "Obrigado!", 0.5, 5.05, 12.3, 1.0, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "ForwardService — \"A Ford saiu da fábrica, não da sua vida.\"", 0.5, 6.1, 12.3, 0.55, size=16, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Código: github.com/fwd-ford   |   Vídeo: [INSERIR LINK]   |   Dúvidas: nos encontre após a apresentação", 0.5, 6.75, 12.3, 0.45, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Salvar ────────────────────────────────────────────────────────────────────
out = r"C:\Users\USER\Desktop\espm\forward-docs\academic\pitch\PITCH.pptx"
prs.save(out)
print(f"PITCH.pptx salvo — {len(prs.slides)} slides")
print()
print("DIVISAO POR INTEGRANTE:")
print("  Slides  1-4  : Joao Victor Franco (RM 556790)")
print("  Slides  5-7  : Lucca Borges       (RM 554608)")
print("  Slides  8-11 : Ruan Melo          (RM 557599)")
print("  Slides 12-14 : Rodrigo Jimenez    (RM 558148)")
print("  Slides 15-16 : Todos")
